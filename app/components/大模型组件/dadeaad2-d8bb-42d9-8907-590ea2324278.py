# -*- coding: utf-8 -*-
import importlib.util
import pathlib
base_path = pathlib.Path(__file__).parent.parent / "base.py"
spec = importlib.util.spec_from_file_location("base", str(base_path))
base_module = importlib.util.module_from_spec(spec)
spec.loader.exec_module(base_module)

# 导入所需项目
BaseComponent = base_module.BaseComponent
PortDefinition = base_module.PortDefinition
PropertyDefinition = base_module.PropertyDefinition
PropertyType = base_module.PropertyType
ArgumentType = base_module.ArgumentType
ConnectionType = base_module.ConnectionType


class Component(BaseComponent):
    name = "文档内容提取"
    category = "大模型组件"
    description = ""
    requirements = "pdfplumber,pandas,python-docx,pptx"

    inputs = [
        PortDefinition(name="file_path", label="文档路径", type=ArgumentType.UPLOAD),
    ]
    outputs = [
        PortDefinition(name="text", label="提取的文本", type=ArgumentType.JSON),
        PortDefinition(name="metadata", label="文档元信息", type=ArgumentType.JSON),
    ]

    properties = {
        "extract_images": PropertyDefinition(
            type=PropertyType.BOOL,
            default=False,
            label="是否提取图片（暂不支持）",
        ),
        "page_range": PropertyDefinition(
            type=PropertyType.TEXT,
            default="",
            label="页码范围（如: 1-5, 留空为全部）",
        ),
    }

    def run(self, params, inputs = None):
        import os
        from pathlib import Path    

        file_path = inputs.file_path if inputs else None
        if not file_path or not os.path.exists(file_path):
            error_msg = "错误：文档路径无效或文件不存在"
            self.logger.error(error_msg)
            return {
                "text": error_msg,
                "metadata": {"error": "File not found"}
            }

        extract_images = params.extract_images
        page_range_str = params.page_range.strip()

        # 解析页码范围（简单支持 "1-3" 或 "2"）
        page_range = None
        if page_range_str:
            try:
                if "-" in page_range_str:
                    start, end = map(int, page_range_str.split("-"))
                    page_range = (start - 1, end)  # 转为 0-based
                else:
                    page_num = int(page_range_str)
                    page_range = (page_num - 1, page_num)
            except Exception as e:
                self.logger.warning(f"页码格式错误，将解析全部页面: {e}")

        file_ext = Path(file_path).suffix.lower()
        text = ""
        metadata = {"file_path": file_path, "format": file_ext, "pages": 0}

        try:
            if file_ext == ".pdf":
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    metadata["pages"] = len(pdf.pages)
                    pages_to_extract = range(len(pdf.pages))
                    if page_range:
                        pages_to_extract = range(
                            max(0, page_range[0]),
                            min(len(pdf.pages), page_range[1])
                        )
                    for i in pages_to_extract:
                        page = pdf.pages[i]
                        text += (page.extract_text() or "") + "\n\n"

            elif file_ext in [".docx", ".doc"]:
                from docx import Document
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + "\n"
                # 表格支持（可选）
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            text += cell.text + " | "
                        text += "\n"

            elif file_ext in [".xlsx", ".xls"]:
                import pandas as pd
                xls = pd.ExcelFile(file_path)
                metadata["sheets"] = xls.sheet_names
                for sheet in xls.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet, dtype=str)
                    text += f"=== Sheet: {sheet} ===\n"
                    text += df.to_string(index=False, na_rep="") + "\n\n"

            elif file_ext in [".pptx", ".ppt"]:
                from pptx import Presentation
                prs = Presentation(file_path)
                metadata["slides"] = len(prs.slides)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"

            elif file_ext == ".txt":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()

            else:
                # 尝试以文本方式读取
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()[:10000]  # 限制长度
                metadata["warning"] = "不支持的格式，尝试以文本读取"

            # 清理多余空白
            text = "\n".join(line.strip() for line in text.splitlines() if line.strip())

            return {
                "text": {"doc": text},
                "metadata": metadata
            }

        except Exception as e:
            raise e
